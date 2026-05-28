# =========================================================
# AI IMPLEMENTATION ENABLEMENT ASSISTANT
# FINAL STABLE MVP
# Streamlit + Gemini + FAISS + TFIDF
# =========================================================

import streamlit as st
import google.generativeai as genai
import tempfile
import os
import uuid
import numpy as np
import faiss
import fitz

from sklearn.feature_extraction.text import (
    TfidfVectorizer
)

from docx import Document
from pptx import Presentation

from langchain_text_splitters import (
    RecursiveCharacterTextSplitter
)

# =========================================================
# PAGE CONFIG
# =========================================================

st.set_page_config(
    page_title="AI Implementation Enablement Assistant",
    layout="wide"
)

st.title(
    "AI Implementation Enablement Assistant"
)

st.caption(
    "Generate Functional, Technical & Support outputs from implementation artifacts"
)

# =========================================================
# CONFIG
# =========================================================

MAX_CONTEXT_LENGTH = 12000
TOP_K_RESULTS = 3

# =========================================================
# GEMINI CONFIG
# =========================================================

st.sidebar.header("Configuration")

gemini_api_key = st.sidebar.text_input(
    "Gemini API Key",
    type="password"
)

if gemini_api_key:

    genai.configure(
        api_key=gemini_api_key
    )

# =========================================================
# IMAGE STORAGE
# =========================================================

IMAGE_DIR = "extracted_images"

os.makedirs(
    IMAGE_DIR,
    exist_ok=True
)

# =========================================================
# SESSION STATE
# =========================================================

if "vector_store" not in st.session_state:
    st.session_state.vector_store = None

if "chunks" not in st.session_state:
    st.session_state.chunks = []

if "metadata" not in st.session_state:
    st.session_state.metadata = []

if "processed" not in st.session_state:
    st.session_state.processed = False

if "vectorizer" not in st.session_state:
    st.session_state.vectorizer = None

# =========================================================
# PDF EXTRACTION
# =========================================================

def extract_pdf_content(
    file_path,
    file_name
):

    pages_data = []

    doc = fitz.open(file_path)

    for page in doc:

        page_text = page.get_text()

        page_images = []

        image_list = page.get_images(
            full=True
        )

        for img in image_list:

            try:

                xref = img[0]

                base_image = (
                    doc.extract_image(xref)
                )

                image_bytes = (
                    base_image["image"]
                )

                image_ext = (
                    base_image["ext"]
                )

                image_name = (
                    f"{uuid.uuid4()}.{image_ext}"
                )

                image_path = os.path.join(
                    IMAGE_DIR,
                    image_name
                )

                with open(
                    image_path,
                    "wb"
                ) as f:

                    f.write(image_bytes)

                page_images.append(
                    image_path
                )

            except:
                pass

        pages_data.append({

            "text": page_text,

            "images": page_images,

            "source_file": file_name
        })

    return pages_data

# =========================================================
# DOCX EXTRACTION
# =========================================================

def extract_docx_content(
    file_path,
    file_name
):

    sections = []

    doc = Document(file_path)

    text = "\n".join(
        [
            para.text
            for para in doc.paragraphs
        ]
    )

    sections.append({

        "text": text,

        "images": [],

        "source_file": file_name
    })

    return sections

# =========================================================
# PPTX EXTRACTION
# =========================================================

def extract_pptx_content(
    file_path,
    file_name
):

    slides_data = []

    prs = Presentation(file_path)

    for slide in prs.slides:

        slide_text = ""

        slide_images = []

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                slide_text += (
                    shape.text + "\n"
                )

            if shape.shape_type == 13:

                try:

                    image = shape.image

                    image_bytes = (
                        image.blob
                    )

                    image_ext = (
                        image.ext
                    )

                    image_name = (
                        f"{uuid.uuid4()}.{image_ext}"
                    )

                    image_path = os.path.join(
                        IMAGE_DIR,
                        image_name
                    )

                    with open(
                        image_path,
                        "wb"
                    ) as f:

                        f.write(image_bytes)

                    slide_images.append(
                        image_path
                    )

                except:
                    pass

        slides_data.append({

            "text": slide_text,

            "images": slide_images,

            "source_file": file_name
        })

    return slides_data

# =========================================================
# TXT EXTRACTION
# =========================================================

def extract_txt_content(
    file_path,
    file_name
):

    with open(
        file_path,
        "r",
        encoding="utf-8"
    ) as f:

        text = f.read()

    return [{

        "text": text,

        "images": [],

        "source_file": file_name
    }]

# =========================================================
# GENERIC EXTRACTION
# =========================================================

def extract_content(uploaded_file):

    suffix = (
        uploaded_file.name
        .split(".")[-1]
        .lower()
    )

    with tempfile.NamedTemporaryFile(
        delete=False,
        suffix=f".{suffix}"
    ) as tmp:

        tmp.write(
            uploaded_file.read()
        )

        temp_path = tmp.name

    try:

        if suffix == "pdf":

            data = extract_pdf_content(
                temp_path,
                uploaded_file.name
            )

        elif suffix == "docx":

            data = extract_docx_content(
                temp_path,
                uploaded_file.name
            )

        elif suffix == "pptx":

            data = extract_pptx_content(
                temp_path,
                uploaded_file.name
            )

        elif suffix == "txt":

            data = extract_txt_content(
                temp_path,
                uploaded_file.name
            )

        else:

            data = []

    finally:

        os.unlink(temp_path)

    return data

# =========================================================
# CHUNKING
# =========================================================

def chunk_text(text):

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=700,
        chunk_overlap=100
    )

    return splitter.split_text(text)

# =========================================================
# VECTOR STORE
# =========================================================

def create_vector_store(chunks):

    vectorizer = TfidfVectorizer()

    embeddings = vectorizer.fit_transform(
        chunks
    ).toarray()

    embeddings = np.array(
        embeddings
    ).astype("float32")

    dimension = embeddings.shape[1]

    index = faiss.IndexFlatL2(
        dimension
    )

    index.add(embeddings)

    return index, vectorizer

# =========================================================
# RETRIEVAL
# =========================================================

def retrieve_chunks(
    query,
    top_k=TOP_K_RESULTS
):

    query_embedding = (
        st.session_state.vectorizer
        .transform([query])
        .toarray()
    )

    query_embedding = np.array(
        query_embedding
    ).astype("float32")

    distances, indices = (
        st.session_state.vector_store.search(
            query_embedding,
            top_k
        )
    )

    retrieved = []

    for idx in indices[0]:

        retrieved.append({

            "chunk":
                st.session_state.chunks[idx],

            "metadata":
                st.session_state.metadata[idx]
        })

    return retrieved

# =========================================================
# GENERATE OUTPUT
# =========================================================

ROLE_RETRIEVAL_QUERIES = {

    "Functional Team": """
    business workflow
    process flow
    approval process
    user actions
    validations
    configurations
    business rules
    end user impact
    """,

    "Technical Team": """
    APIs
    backend logic
    database changes
    integrations
    deployment
    scheduler
    dependencies
    technical validations
    performance
    """,

    "Support Team": """
    troubleshooting
    errors
    failures
    support cases
    logs
    retry mechanism
    validation failures
    escalation
    issue resolution
    """
}

def generate_output(
    team,
    additional_instruction
):

    retrieval_query = ROLE_RETRIEVAL_QUERIES[team]

    retrieval_query += additional_instruction

    retrieved_chunks = retrieve_chunks(
        retrieval_query
    )

    if len(retrieved_chunks) == 0:

        return (
            "No relevant content found.",
            []
        )

    context = ""

    related_images = []

    for item in retrieved_chunks:

        context += f"""

        FILE:
        {item['metadata']['file_name']}

        CONTENT:
        {item['chunk']}

        ------------------------------
        """

        for img in item["metadata"]["images"]:

            if os.path.exists(img):

                related_images.append(img)

    related_images = list(
        set(related_images)
    )[:5]

    context = context[
        :MAX_CONTEXT_LENGTH
    ]

    if team == "Functional Team":

    prompt = f"""
    You are a Functional Consultant.

    CONTEXT:
    {context}

    TASK:
    Generate a business-oriented implementation summary.

    FORMAT:
    # Business Overview
    # User Workflow
    # Configurations
    # Validation Rules
    # Business Impact
    # Key User Actions

    STYLE:
    - Business friendly
    - Process oriented
    - Functional language only
    - Avoid deep technical jargon

    ADDITIONAL INSTRUCTIONS:
    {additional_instruction}
    """

elif team == "Technical Team":

    prompt = f"""
    You are a Technical Architect.

    CONTEXT:
    {context}

    TASK:
    Generate a technical implementation summary.

    FORMAT:
    # Technical Overview
    # APIs Impacted
    # Database Changes
    # Integrations
    # Deployment Impact
    # Dependencies
    # Technical Risks

    STYLE:
    - Engineering focused
    - Technical terminology
    - Implementation oriented
    - Mention system impacts

    ADDITIONAL INSTRUCTIONS:
    {additional_instruction}
    """

else:

    prompt = f"""
    You are a Support Engineer.

    CONTEXT:
    {context}

    TASK:
    Generate a support troubleshooting guide.

    FORMAT:
    # Common Issues
    # Error Scenarios
    # Troubleshooting Steps
    # Logs to Check
    # Retry Mechanism
    # Escalation Guidance

    STYLE:
    - Operational
    - Troubleshooting oriented
    - Focus on issue handling
    - Mention validations and failures

    ADDITIONAL INSTRUCTIONS:
    {additional_instruction}
    """
    try:

        model = genai.GenerativeModel(
            "gemini-2.5-flash"
        )

        response = model.generate_content(
            prompt
        )

        output_text = response.text

    except Exception as e:

        output_text = (
            f"Generation Error: {str(e)}"
        )

    return (
        output_text,
        related_images
    )
# =========================================================
# UPLOAD SCREEN
# =========================================================

st.header(
    "1. Upload Implementation Artifacts"
)

uploaded_files = st.file_uploader(
    "Upload PDF / DOCX / PPTX / TXT",
    type=["pdf", "docx", "pptx", "txt"],
    accept_multiple_files=True
)

# =========================================================
# PROCESS FILES
# =========================================================

if uploaded_files:

    if st.button("Process Files"):

        if not gemini_api_key:

            st.error(
                "Please enter Gemini API Key"
            )

            st.stop()

        all_chunks = []
        all_metadata = []

        st.subheader(
            "Processing Files..."
        )

        for file in uploaded_files:

            st.write(
                f"Reading: {file.name}"
            )

            extracted_sections = (
                extract_content(file)
            )

            for section in extracted_sections:

                text = section["text"]

                images = section["images"]

                if not text.strip():

                    continue

                chunks = chunk_text(text)

                for idx, chunk in enumerate(
                    chunks
                ):

                    all_chunks.append(chunk)

                    all_metadata.append({

                        "file_name":
                            section["source_file"],

                        "chunk_id":
                            idx,

                        "images":
                            images
                    })

        if len(all_chunks) == 0:

            st.error(
                "No readable content found"
            )

            st.stop()

        st.subheader(
            "Creating Knowledge Index..."
        )

        vector_store, vectorizer = (
            create_vector_store(
                all_chunks
            )
        )

        st.session_state.vector_store = (
            vector_store
        )

        st.session_state.vectorizer = (
            vectorizer
        )

        st.session_state.chunks = (
            all_chunks
        )

        st.session_state.metadata = (
            all_metadata
        )

        st.session_state.processed = True

        st.success(
            "Knowledge Processing Completed"
        )

# =========================================================
# OUTPUT SCREEN
# =========================================================

if st.session_state.processed:

    st.header(
        "2. Generate Team-Specific Output"
    )

    team = st.selectbox(
        "Generate For",
        [
            "Functional Team",
            "Technical Team",
            "Support Team"
        ]
    )

    additional_instruction = st.text_area(
        "Additional Instructions",
        placeholder="""
Examples:
- Focus on approval workflow
- Explain integrations
- Include troubleshooting
        """
    )

    if st.button(
        "Generate Output"
    ):

        with st.spinner(
            "Generating Output..."
        ):

            output, images = (
                generate_output(
                    team,
                    additional_instruction
                )
            )

        st.subheader(
            "Generated Output"
        )

        st.markdown(output)

        if images:

            st.subheader(
                "Relevant Screenshots"
            )

            for img in images:

                st.image(
                    img,
                    use_container_width=True
                )

# =========================================================
# FOOTER
# =========================================================

st.divider()

st.caption(
    "Final Stable MVP"
)
